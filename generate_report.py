"""MazicAlign 프로젝트 분석 보고서 생성기 - PPTX + XLSX"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

# ─────────────────────────── 색상 팔레트 ───────────────────────────
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD    = RGBColor(0x22, 0x22, 0x3B)
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xAA, 0xAA, 0xBB)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.space_after = Pt(6)
    return txBox


def add_card(slide, left, top, width, height, title, body_items,
             title_color=ACCENT, body_color=WHITE):
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    shape.shadow.inherit = False

    # Title
    add_textbox(slide, left + 0.2, top + 0.15, width - 0.4, 0.4,
                title, font_size=16, color=title_color, bold=True)
    # Body
    add_bullet_list(slide, left + 0.2, top + 0.55, width - 0.4, height - 0.7,
                    body_items, font_size=12, color=body_color)


def add_table_slide(slide, left, top, rows_data, col_widths, header_color=ACCENT):
    """rows_data: list of lists. First row = header."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                         Inches(sum(col_widths)), Inches(rows * 0.45))
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "맑은 고딕"
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = WHITE

            # cell fill
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x7C, 0x3A, 0xED)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD if ri % 2 == 1 else RGBColor(0x2D, 0x2D, 0x48)

    return table_shape


# ═══════════════════════════════════════════════════════════════
#                    PPTX 생성
# ═══════════════════════════════════════════════════════════════
def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    # ── Slide 1: Title ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)

    # Accent bar
    shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    add_textbox(sl, 1.5, 1.8, 10, 1, "MazicAlign", font_size=52, color=ACCENT, bold=True)
    add_textbox(sl, 1.5, 2.8, 10, 0.8, "로컬 전용 3D STL 뷰어 & 슬라이서 프로젝트 분석 보고서",
                font_size=24, color=WHITE)
    add_textbox(sl, 1.5, 4.0, 10, 0.6,
                "프로젝트 유형: 치과/교정 전문 3D 프린팅 워크플로우 웹 애플리케이션",
                font_size=16, color=GRAY)
    add_textbox(sl, 1.5, 4.6, 10, 0.4,
                "분석 일자: 2026-03-17  |  총 코드량: ~5,545 Lines  |  TypeScript + React + BabylonJS",
                font_size=14, color=GRAY)

    # ── Slide 2: 프로젝트 개요 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "1. 프로젝트 개요", font_size=32, color=ACCENT, bold=True)

    add_card(sl, 0.5, 1.2, 5.8, 2.8, "프로젝트 목적", [
        "치과/교정 분야 3D STL 모델 뷰어 & 슬라이서",
        "로컬 PC(Windows 10/11) 전용 - 인터넷 불필요",
        "ZIP 배포 → install.bat → start.bat 실행",
        "단일 Express 서버 (포트 5173)로 구동",
        "SQLite 단일 파일 DB로 데이터 관리",
    ])

    add_card(sl, 7.0, 1.2, 5.8, 2.8, "핵심 특징", [
        "브라우저 기반 3D 뷰어 (WebGL / BabylonJS)",
        "FDM + DLP 듀얼 슬라이싱 엔진 내장",
        "Web Worker로 비차단(non-blocking) 슬라이싱",
        "로컬 파일시스템 직접 탐색 & 임포트",
        "변환 이력(History) 자동 기록 & 조회",
    ])

    add_card(sl, 0.5, 4.3, 12.3, 2.8, "배포 & 실행 방식", [
        "1) 프로젝트 전체를 ZIP으로 압축 (node_modules 제외)",
        "2) 사용자가 ZIP 압축 해제 후 install.bat 실행 → npm 패키지 설치 (~3분)",
        "3) build.bat 실행 → 프론트엔드(Vite) + 백엔드(tsc) 빌드",
        "4) start.bat 실행 → http://localhost:5173 자동 오픈",
        "5) 이후 인터넷 불필요, 모든 데이터 로컬 저장 (SQLite + 파일시스템)",
    ])

    # ── Slide 3: 기술 스택 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "2. 기술 스택 (Tech Stack)", font_size=32, color=ACCENT, bold=True)

    fe_data = [
        ["카테고리", "기술", "버전", "용도"],
        ["프레임워크", "React", "18.2.0", "컴포넌트 UI"],
        ["언어", "TypeScript", "5.2.2", "타입 안전성"],
        ["빌드 도구", "Vite", "5.0.8", "번들링 & HMR"],
        ["3D 엔진", "BabylonJS", "6.36.0", "WebGL 렌더링"],
        ["스타일링", "Tailwind CSS", "3.3.6", "유틸리티 CSS"],
        ["상태 관리", "Zustand", "4.4.7", "경량 상태 관리"],
        ["라우팅", "React Router", "6.20.0", "SPA 라우팅"],
        ["폴리곤 연산", "js-clipper", "1.0.1", "폴리곤 Boolean 연산"],
    ]
    add_textbox(sl, 0.5, 1.0, 3, 0.4, "Frontend", font_size=18, color=GREEN, bold=True)
    add_table_slide(sl, 0.5, 1.5, fe_data, [1.5, 1.8, 1.0, 2.0])

    be_data = [
        ["카테고리", "기술", "버전", "용도"],
        ["런타임", "Node.js", "22 LTS", "서버 런타임"],
        ["프레임워크", "Express", "4.18.2", "HTTP API 서버"],
        ["데이터베이스", "SQLite (better-sqlite3)", "11.0.0", "단일 파일 RDB"],
        ["파일 업로드", "Multer", "1.4.5", "멀티파트 처리"],
        ["ID 생성", "UUID", "9.0.1", "고유 식별자"],
        ["환경설정", "dotenv", "16.6.1", "환경변수 관리"],
    ]
    add_textbox(sl, 7.0, 1.0, 3, 0.4, "Backend", font_size=18, color=GREEN, bold=True)
    add_table_slide(sl, 7.0, 1.5, be_data, [1.5, 1.8, 1.0, 2.0])

    # ── Slide 4: 시스템 아키텍처 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "3. 시스템 아키텍처", font_size=32, color=ACCENT, bold=True)

    # Architecture diagram (text-based)
    arch_text = """
┌─────────────────────────────────────────────────────────────────────────┐
│                           브라우저 (localhost:5173)                       │
│  ┌──────────────┐  ┌──────────────────┐  ┌────────────────────────┐    │
│  │  ProjectPage  │  │   ViewerPage      │  │   Slicer (Web Worker)  │    │
│  │  프로젝트 관리 │  │   3D 뷰어          │  │   FDM G-code 생성      │    │
│  │               │  │   BabylonJS        │  │   DLP 마스크 이미지     │    │
│  └──────────────┘  └──────────────────┘  └────────────────────────┘    │
│                              │  REST API (fetch)                         │
└──────────────────────────────┼───────────────────────────────────────────┘
                               ▼
┌──────────────────────────────────────────────────────────────────────────┐
│                        Express 서버 (Node.js)                             │
│  ┌─────────────┐  ┌──────────────┐  ┌─────────────┐  ┌────────────┐   │
│  │  /api/projects│  │  /api/stl     │  │  /api/fs     │  │  /uploads   │   │
│  │  프로젝트 CRUD │  │  STL파일 CRUD  │  │  파일시스템   │  │  정적 파일   │   │
│  └──────┬──────┘  └──────┬───────┘  └──────┬──────┘  └────────────┘   │
│         │                │                  │                            │
│         ▼                ▼                  ▼                            │
│  ┌─────────────────────────────┐  ┌────────────────────┐               │
│  │  SQLite (mazicalign.db)      │  │  Windows 파일시스템  │               │
│  │  projects / stl_files /      │  │  C:\\ D:\\ ...        │               │
│  │  adjustment_logs / users     │  │  STL 파일 직접 접근   │               │
│  └─────────────────────────────┘  └────────────────────┘               │
└──────────────────────────────────────────────────────────────────────────┘
"""
    txBox = add_textbox(sl, 0.3, 1.1, 12.7, 6.2, arch_text.strip(),
                        font_size=11, color=WHITE, font_name="Consolas")
    txBox.text_frame.paragraphs[0].font.name = "Consolas"
    for p in txBox.text_frame.paragraphs:
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

    # ── Slide 5: 핵심 기능 상세 (1) ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "4. 핵심 기능 상세 (1/2)", font_size=32, color=ACCENT, bold=True)

    add_card(sl, 0.5, 1.2, 6.0, 2.5, "3D 뷰어 (BabylonJS)", [
        "ASCII & Binary STL 파일 동시 로딩",
        "ArcRotateCamera (궤도/팬/줌)",
        "Gizmo Manager: 이동/회전/크기 조절 드래그",
        "파일별 가시성(Visibility) 토글 & 투명도(0~100%)",
        "선택 모델 하이라이트 & 커스텀 머티리얼",
    ])

    add_card(sl, 7.0, 1.2, 6.0, 2.5, "프로젝트 관리", [
        "프로젝트 생성/조회/수정/삭제 (CRUD)",
        "8자 랜덤 프로젝트 코드 자동 생성",
        "환자 정보(patientInfo) JSON 저장",
        "프로젝트별 STL 파일 분리 관리",
        "SQLite 기반 영속 저장",
    ])

    add_card(sl, 0.5, 4.0, 6.0, 3.0, "변환 컨트롤 & 이력", [
        "Translation (X/Y/Z) - 수치 입력 또는 Gizmo 드래그",
        "Rotation (X/Y/Z 오일러각) - 수치 입력 또는 Gizmo",
        "Scale (균등 또는 축별) - 수치 입력",
        "모든 변환 자동 기록 (AdjustmentLog)",
        "이력 조회: STL 파일별, 최신순 정렬",
        "개별 로그 삭제 또는 전체 초기화 가능",
    ])

    add_card(sl, 7.0, 4.0, 6.0, 3.0, "로컬 파일 탐색기", [
        "백엔드 /api/fs 엔드포인트로 PC 파일시스템 직접 탐색",
        "Windows 드라이브 문자(A:~Z:) 지원",
        "디렉토리 + .stl 파일만 필터링 표시",
        "파일 메타데이터(크기, 타입) 반환",
        "선택한 STL 파일 즉시 프로젝트로 임포트",
        "파일 업로드(Multer)와 로컬 경로 임포트 모두 지원",
    ])

    # ── Slide 6: 핵심 기능 상세 (2) - 슬라이서 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "5. 핵심 기능 상세 (2/2) - 슬라이서", font_size=32, color=ACCENT, bold=True)

    add_card(sl, 0.5, 1.2, 6.0, 5.5, "FDM 슬라이서 (G-code 생성)", [
        "[프로세스]",
        "1. 3D 메시의 Z축 범위 계산",
        "2. 레이어 높이 간격으로 Z-plane 반복",
        "3. 삼각형-평면 교차점 계산 → 윤곽선 생성",
        "4. 윤곽선 오프셋 → 동심원 벽(Walls) 생성",
        "5. 내부 영역 인필(Infill) 패턴 채우기",
        "6. G-code 명령어(G0/G1 + E값) 출력",
        "",
        "[설정 파라미터]",
        "• 레이어 높이: 0.05 ~ 0.5mm",
        "• 노즐 직경 (mm)",
        "• 벽 개수 (shells), 인필 밀도 (%)",
        "• 인필 패턴: lines / grid / zigzag",
        "• 프린트 속도 (mm/s)",
        "• 벽 순서: inner→outer / outer→inner",
        "• 출력 순서: walls-first / infill-first",
    ])

    add_card(sl, 7.0, 1.2, 6.0, 5.5, "DLP 슬라이서 (마스크 이미지 생성)", [
        "[프로세스]",
        "1. 3D 메시를 레이어별 2D 폴리곤으로 슬라이싱",
        "2. OffscreenCanvas에 폴리곤 렌더링",
        "   (흰색 = 노출, 검정 = 비노출)",
        "3. Canvas → PNG Blob 변환",
        "4. Data URL 또는 이미지 시퀀스로 내보내기",
        "",
        "[설정 파라미터]",
        "• 레이어 높이: 0.01 ~ 0.1mm",
        "• 빌드 플레이트 크기 (너비 × 깊이 mm)",
        "• 프로젝터 해상도 (X × Y 픽셀)",
        "• 픽셀 크기 (마이크론) - 자동 계산 가능",
        "• 광원 파워 (%)",
        "• 노출 시간 (초)",
        "• Z-리프트 속도 (mm/s)",
        "",
        "[핵심 라이브러리]",
        "• js-clipper: 폴리곤 오프셋 & Boolean 연산",
        "• OffscreenCanvas: Web Worker 내 렌더링",
    ])

    # ── Slide 7: 디렉토리 구조 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "6. 프로젝트 디렉토리 구조", font_size=32, color=ACCENT, bold=True)

    dir_text = """MazicAlign/
├── frontend/                   # React 18 + Vite SPA (~3,200 lines)
│   ├── src/
│   │   ├── pages/              # ProjectPage, ViewerPage, HomePage, LoginPage
│   │   ├── components/         # STLViewer, STLFileList, LocalFileBrowser,
│   │   │   └── Slicer/         # SlicerPanel, SlicePreview
│   │   ├── hooks/              # useAuth, useProjects, useSTLFiles, useAdjustmentLogs
│   │   ├── services/           # auth, project, stl API 서비스
│   │   │   └── slicer/         # SliceEngine, GCodeGenerator, ImageGenerator,
│   │   │                       # PolygonClipper, SlicerService, SlicerWorker
│   │   ├── utils/              # babylon.utils, stl-loader.utils, transform.utils
│   │   ├── types/              # stl, project, api, auth 타입 정의
│   │   └── config/             # firebase.config (레거시, 미사용)
│   ├── vite.config.ts          # Vite 설정 + API 프록시
│   └── package.json
│
├── backend/                    # Express API 서버 (~900 lines)
│   ├── src/
│   │   ├── controllers/        # project, stl HTTP 핸들러
│   │   ├── services/           # project, stl 비즈니스 로직
│   │   ├── routes/             # project, stl, fs 라우터
│   │   ├── models/             # project, stl, user 타입
│   │   └── config/             # database (SQLite), server.config
│   ├── data/                   # mazicalign.db (자동 생성)
│   └── uploads/stl/            # STL 파일 저장소
│
├── install.bat / build.bat / start.bat   # Windows 배치 스크립트
├── README.md / DESIGN.md / SETUP.md      # 문서
└── cors.json                              # CORS 설정"""

    txBox = add_textbox(sl, 0.5, 1.1, 12.3, 6.2, dir_text, font_size=12, color=WHITE, font_name="Consolas")
    for p in txBox.text_frame.paragraphs:
        p.font.name = "Consolas"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

    # ── Slide 8: API 엔드포인트 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "7. REST API 엔드포인트", font_size=32, color=ACCENT, bold=True)

    api_data = [
        ["메서드", "엔드포인트", "설명"],
        ["GET", "/health", "서버 상태 확인"],
        ["GET", "/api/fs?path=...", "로컬 파일시스템 탐색"],
        ["GET", "/api/projects?ownerId=...", "프로젝트 목록 조회"],
        ["POST", "/api/projects", "프로젝트 생성"],
        ["PUT", "/api/projects/{id}", "프로젝트 수정"],
        ["DELETE", "/api/projects/{id}", "프로젝트 삭제"],
        ["GET", "/api/stl?projectId=...", "STL 파일 목록"],
        ["POST", "/api/stl/upload", "STL 파일 업로드 (멀티파트)"],
        ["POST", "/api/stl/import-path", "로컬 경로에서 STL 임포트"],
        ["PUT", "/api/stl/{id}/visibility", "가시성 토글"],
        ["PUT", "/api/stl/{id}/transform", "변환값 업데이트"],
        ["DELETE", "/api/stl/{id}", "STL 파일 삭제"],
        ["GET", "/api/stl/{id}/logs", "조정 이력 조회"],
        ["POST", "/api/stl/{id}/logs", "조정 이력 생성"],
        ["DELETE", "/api/stl/{id}/logs/{logId}", "개별 이력 삭제"],
        ["DELETE", "/api/stl/{id}/logs", "전체 이력 초기화"],
    ]
    add_table_slide(sl, 0.5, 1.1, api_data, [1.0, 3.5, 3.0])

    # ── Slide 9: 데이터 모델 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "8. 데이터 모델 (SQLite)", font_size=32, color=ACCENT, bold=True)

    db_data = [
        ["테이블", "주요 필드", "설명"],
        ["users", "userId, email, displayName, role, createdAt", "사용자 (로컬 단일 사용자)"],
        ["projects", "projectId, ownerId, projectCode, projectName, patientInfo", "프로젝트 (환자 정보 포함)"],
        ["stl_files", "stlId, projectId, originalUrl, fileName, fileSize, visibility, currentTransform", "STL 파일 메타데이터 + 변환값"],
        ["adjustment_logs", "logId, projectId, stlId, userId, adjustmentType, deltaValue, timestamp", "변환 이력 로그"],
    ]
    add_table_slide(sl, 0.5, 1.2, db_data, [2.0, 5.5, 4.0])

    add_card(sl, 0.5, 3.8, 12.3, 3.2, "Transform 데이터 구조 (JSON)", [
        "currentTransform = {",
        "  translation: { x: number, y: number, z: number }",
        "  rotation: { x: number, y: number, z: number, w: number }  // Quaternion",
        "  scale: { x: number, y: number, z: number }",
        "}",
        "",
        "AdjustmentLog.deltaValue = Partial<Vector3> | Partial<Quaternion>",
        "AdjustmentLog.adjustmentType = 'Translation' | 'Rotation' | 'Scale'",
    ])

    # ── Slide 10: 참고/영향받은 기술 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "9. 참고한 기술 & 영향 분석", font_size=32, color=ACCENT, bold=True)

    add_card(sl, 0.5, 1.2, 6.0, 2.8, "3D 뷰어 참고", [
        "BabylonJS 공식 문서 - STL Loader, Gizmo Manager",
        "BabylonJS Playground 예제 - ArcRotateCamera 설정",
        "3D 모델 뷰어 패턴 (Blender, Meshmixer UI 참고 가능성)",
        "WebGL 기반 치과 뷰어 상용 솔루션 참고 가능성",
        "  (예: 3Shape, exocad 등의 UI/UX 패턴)",
    ])

    add_card(sl, 7.0, 1.2, 6.0, 2.8, "슬라이서 참고", [
        "오픈소스 슬라이서: Cura, PrusaSlicer 알고리즘 참고",
        "Z-plane 교차 알고리즘 - 일반적인 메시 슬라이싱 기법",
        "js-clipper: Angus Johnson의 Clipper 라이브러리 JS 포팅",
        "G-code 표준: RepRap G-code 명령어 체계",
        "DLP 슬라이싱: ChiTuBox, Lychee 등 DLP 슬라이서 참고",
    ])

    add_card(sl, 0.5, 4.3, 6.0, 2.8, "아키텍처 참고", [
        "Create React App → Vite 마이그레이션 패턴",
        "Express + SQLite 로컬 서버 패턴",
        "Zustand 상태 관리 패턴 (공식 문서)",
        "Multer 파일 업로드 패턴 (Express 공식 예제)",
        "Tailwind CSS 유틸리티 패턴",
    ])

    add_card(sl, 7.0, 4.3, 6.0, 2.8, "V2 로드맵 (DESIGN.md)", [
        "BabylonJS → Three.js + @react-three/fiber 전환 계획",
        "  → 번들 크기 감소 (Three.js ≈ BabylonJS의 절반)",
        "Gizmo → @react-three/drei TransformControls",
        "아이콘: Lucide React 도입 예정",
        "슬라이서: Web Worker 방식 유지, 성능 개선",
    ])

    # ── Slide 11: PrusaSlicer 기능 비교표 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6,
                "10. MazicAlign vs PrusaSlicer 기능 비교 (FDM)",
                font_size=32, color=ACCENT, bold=True)
    add_textbox(sl, 0.5, 0.85, 12, 0.3,
                "참고: github.com/JoWooHyun/PrusaSlicer  |  C++ ~68.9%  |  GNU AGPL v3  |  Slic3r 기반",
                font_size=12, color=GRAY)

    RED = RGBColor(0xEF, 0x44, 0x44)
    cmp_data = [
        ["기능 영역", "기능", "PrusaSlicer", "MazicAlign", "Gap"],
        ["인필 패턴", "Lines / Grid / Zigzag", "O", "O", "-"],
        ["인필 패턴", "Honeycomb (벌집)", "O", "X", "추가 필요"],
        ["인필 패턴", "3D Honeycomb", "O", "X", "추가 필요"],
        ["인필 패턴", "Gyroid (자이로이드)", "O", "X", "추가 필요"],
        ["인필 패턴", "Concentric (동심원)", "O", "X", "추가 필요"],
        ["인필 패턴", "Lightning (번개)", "O", "X", "추가 필요"],
        ["인필 패턴", "Adaptive (적응형)", "O", "X", "추가 필요"],
        ["벽 생성", "기본 벽(Walls/Perimeters)", "O", "O", "-"],
        ["벽 생성", "Arachne 가변폭 윤곽선", "O", "X", "고급 기능"],
        ["벽 생성", "Seam 배치 최적화 (10+전략)", "O", "X", "추가 필요"],
        ["서포트", "일반 서포트 생성", "O", "X", "추가 필요"],
        ["서포트", "트리 서포트 (Tree)", "O", "X", "추가 필요"],
        ["서포트", "유기적 서포트 (Organic)", "O", "X", "추가 필요"],
        ["G-code", "기본 G0/G1 생성", "O", "O", "-"],
        ["G-code", "다중 펌웨어 (Marlin/Prusa/RepRap)", "O", "X", "추가 필요"],
        ["G-code", "리트랙션 & 압력 균등화", "O", "X", "추가 필요"],
        ["G-code", "냉각 제어 (CoolingBuffer)", "O", "X", "추가 필요"],
        ["G-code", "Wipe Tower (다중 재료)", "O", "X", "해당없음"],
        ["G-code", "나선형 화병 모드 (Spiral Vase)", "O", "X", "추가 필요"],
        ["G-code", "이동 경로 최적화 (Travels)", "O", "X", "추가 필요"],
        ["G-code", "경계선 교차 회피", "O", "X", "추가 필요"],
        ["메시 처리", "STL 자동 복구", "O", "X", "추가 필요"],
        ["메시 처리", "다중 포맷 (STL/OBJ/AMF/3MF)", "O", "△ (STL만)", "추가 권장"],
        ["메시 처리", "CSG Boolean 연산", "O", "X", "고급 기능"],
        ["후처리", "후처리 스크립트", "O", "X", "추가 권장"],
        ["후처리", "G-code 매크로 커스터마이징", "O", "X", "추가 권장"],
        ["후처리", "썸네일 생성", "O", "X", "추가 권장"],
        ["SLA/DLP", "마스크 이미지 생성", "O (mSLA)", "O", "-"],
        ["SLA/DLP", "자동 서포트 포인트 생성", "O", "X", "추가 필요"],
        ["SLA/DLP", "모델 중공화 (Hollowing)", "O", "X", "추가 필요"],
        ["SLA/DLP", "빌드 패드 (Pad)", "O", "X", "추가 필요"],
        ["SLA/DLP", "최적 회전 탐색 (Rotfinder)", "O", "X", "추가 권장"],
        ["SLA/DLP", "Z축 보정 (ZCorrection)", "O", "X", "추가 권장"],
        ["기타", "CLI 모드 (GUI 없이)", "O", "X", "해당없음"],
        ["기타", "멀티스레드 처리", "O (C++ Thread)", "O (Web Worker)", "-"],
        ["기타", "자동화 유닛 테스트", "O", "X", "추가 권장"],
    ]
    add_table_slide(sl, 0.3, 1.2, cmp_data[:20], [1.3, 2.8, 1.3, 1.3, 1.3])

    # ── Slide 12: 비교표 (2/2) ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6,
                "10. MazicAlign vs PrusaSlicer 기능 비교 (계속)",
                font_size=32, color=ACCENT, bold=True)

    add_table_slide(sl, 0.3, 1.0, [cmp_data[0]] + cmp_data[20:], [1.3, 2.8, 1.3, 1.3, 1.3])

    # Summary card
    add_card(sl, 7.5, 1.0, 5.5, 6.0, "Gap 분석 요약", [
        "[FDM 슬라이서 주요 부족 기능]",
        "• 인필 패턴: 7종 → 현재 3종 (Honeycomb/Gyroid 등 부재)",
        "• 서포트 생성: 일반/트리/유기적 모두 미구현",
        "• Seam 배치 최적화 (10+전략) 미구현",
        "• 리트랙션/냉각/압력균등화 등 G-code 고급 기능",
        "• 이동 경로 최적화 & 경계선 교차 회피",
        "• STL 자동 복구 & 다중 포맷 지원",
        "• 나선형 화병 모드 (Spiral Vase)",
        "",
        "[DLP/SLA 주요 부족 기능]",
        "• 자동 서포트 포인트 생성",
        "• 모델 중공화 (Hollowing)",
        "• 빌드 패드 (Pad) 생성",
        "• 최적 인쇄 회전 탐색",
        "• Z축 보정",
        "",
        "[MazicAlign 고유 강점]",
        "• 브라우저 기반 (설치 불필요)",
        "• 치과/교정 전문 워크플로우",
        "• 로컬 파일시스템 직접 탐색",
        "• 변환 이력 자동 기록",
    ])

    # ── Slide 13: 코드 통계 & 요약 ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, 0.5, 0.3, 12, 0.6, "11. 코드 통계 & 종합 평가", font_size=32, color=ACCENT, bold=True)

    stats_data = [
        ["구분", "라인 수", "비중", "주요 내용"],
        ["Frontend (페이지/컴포넌트)", "~1,500", "27%", "ViewerPage, STLViewer, SlicerPanel 등"],
        ["Frontend (서비스/슬라이서)", "~1,200", "22%", "SliceEngine, GCodeGenerator, ImageGenerator"],
        ["Frontend (hooks/utils/types)", "~500", "9%", "커스텀 훅, BabylonJS 유틸, 타입 정의"],
        ["Backend (전체)", "~900", "16%", "Express API, SQLite, 파일 관리"],
        ["설정/타입 정의", "~1,445", "26%", "tsconfig, vite.config, 타입 파일"],
        ["합계", "~5,545", "100%", "TypeScript/TSX"],
    ]
    add_table_slide(sl, 0.5, 1.2, stats_data, [2.5, 1.2, 1.0, 4.5])

    add_card(sl, 0.5, 5.0, 12.3, 2.2, "종합 평가", [
        "잘 설계된 풀스택 웹 애플리케이션으로, 3D 프린팅(치과/교정) 워크플로우에 특화되어 있음",
        "프론트엔드/백엔드 분리가 깔끔하며, 서비스 레이어 패턴을 준수함",
        "Web Worker 슬라이싱으로 UI 블로킹 방지 — 성능 고려가 잘 되어 있음",
        "로컬 전용 설계로 보안 우려 최소화, 배포 편의성 극대화",
        "V2에서 Three.js 전환 예정 — 번들 크기 최적화 & React 에코시스템 일관성 확보 목적",
    ])

    prs.save("c:/Users/JoWooHyun/Documents/MazicAlign/MazicAlign_분석보고서_v2.pptx")
    print("PPTX 생성 완료: MazicAlign_분석보고서_v2.pptx")


# ═══════════════════════════════════════════════════════════════
#                    XLSX 생성
# ═══════════════════════════════════════════════════════════════
def create_xlsx():
    wb = openpyxl.Workbook()

    # ── 색상/스타일 정의 ──
    header_fill = PatternFill(start_color="7C3AED", end_color="7C3AED", fill_type="solid")
    header_font = Font(name="맑은 고딕", size=11, bold=True, color="FFFFFF")
    cell_font = Font(name="맑은 고딕", size=10)
    thin_border = Border(
        left=Side(style='thin', color='CCCCCC'),
        right=Side(style='thin', color='CCCCCC'),
        top=Side(style='thin', color='CCCCCC'),
        bottom=Side(style='thin', color='CCCCCC')
    )
    alt_fill_1 = PatternFill(start_color="F5F3FF", end_color="F5F3FF", fill_type="solid")
    alt_fill_2 = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")

    def style_header(ws, row=1, max_col=None):
        for cell in ws[row]:
            if max_col and cell.column > max_col:
                break
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center', vertical='center')
            cell.border = thin_border

    def style_body(ws, start_row=2, max_col=None, max_row=None):
        mr = max_row or ws.max_row
        mc = max_col or ws.max_column
        for ri in range(start_row, mr + 1):
            for ci in range(1, mc + 1):
                cell = ws.cell(row=ri, column=ci)
                cell.font = cell_font
                cell.border = thin_border
                cell.fill = alt_fill_1 if ri % 2 == 0 else alt_fill_2
                cell.alignment = Alignment(vertical='center', wrap_text=True)

    # ── Sheet 1: 기술 스택 ──
    ws = wb.active
    ws.title = "기술 스택"
    headers = ["영역", "카테고리", "기술명", "버전", "용도", "npm 패키지명", "참고 URL"]
    ws.append(headers)

    tech_rows = [
        ["Frontend", "프레임워크", "React", "18.2.0", "컴포넌트 기반 UI", "react", "https://react.dev"],
        ["Frontend", "언어", "TypeScript", "5.2.2", "타입 안전성", "typescript", "https://typescriptlang.org"],
        ["Frontend", "빌드 도구", "Vite", "5.0.8", "번들링 & HMR", "vite", "https://vitejs.dev"],
        ["Frontend", "3D 엔진", "BabylonJS Core", "6.36.0", "WebGL 3D 렌더링", "@babylonjs/core", "https://babylonjs.com"],
        ["Frontend", "3D 로더", "BabylonJS Loaders", "6.36.0", "STL 파일 로딩", "@babylonjs/loaders", "https://babylonjs.com"],
        ["Frontend", "스타일링", "Tailwind CSS", "3.3.6", "유틸리티 CSS 프레임워크", "tailwindcss", "https://tailwindcss.com"],
        ["Frontend", "상태 관리", "Zustand", "4.4.7", "경량 상태 관리", "zustand", "https://zustand-demo.pmnd.rs"],
        ["Frontend", "라우팅", "React Router", "6.20.0", "SPA 라우팅", "react-router-dom", "https://reactrouter.com"],
        ["Frontend", "폴리곤 연산", "js-clipper", "1.0.1", "폴리곤 Boolean/오프셋 연산", "js-clipper", "https://github.com/nicklockwood/js-clipper"],
        ["Frontend", "린팅", "ESLint", "8.55.0", "코드 품질 검사", "eslint", "https://eslint.org"],
        ["Backend", "런타임", "Node.js", "22 LTS", "JavaScript 서버 런타임", "-", "https://nodejs.org"],
        ["Backend", "프레임워크", "Express", "4.18.2", "HTTP API 서버", "express", "https://expressjs.com"],
        ["Backend", "데이터베이스", "SQLite (better-sqlite3)", "11.0.0", "로컬 단일 파일 RDB", "better-sqlite3", "https://github.com/WiseLibs/better-sqlite3"],
        ["Backend", "파일 업로드", "Multer", "1.4.5", "멀티파트 폼 데이터 처리", "multer", "https://github.com/expressjs/multer"],
        ["Backend", "ID 생성", "UUID", "9.0.1", "고유 식별자 생성", "uuid", "https://github.com/uuidjs/uuid"],
        ["Backend", "CORS", "cors", "2.8.5", "Cross-Origin 요청 처리", "cors", "https://github.com/expressjs/cors"],
        ["Backend", "환경설정", "dotenv", "16.6.1", "환경변수 파일 로드", "dotenv", "https://github.com/motdotla/dotenv"],
        ["Backend", "빌드", "tsc-alias", "1.8.16", "TS 경로 별칭 해석", "tsc-alias", "https://github.com/justkey007/tsc-alias"],
        ["Backend", "개발 실행", "tsx", "4.7.0", "TypeScript 직접 실행", "tsx", "https://github.com/privatenumber/tsx"],
    ]
    for row in tech_rows:
        ws.append(row)

    ws.column_dimensions['A'].width = 12
    ws.column_dimensions['B'].width = 14
    ws.column_dimensions['C'].width = 22
    ws.column_dimensions['D'].width = 10
    ws.column_dimensions['E'].width = 28
    ws.column_dimensions['F'].width = 22
    ws.column_dimensions['G'].width = 45
    style_header(ws)
    style_body(ws)

    # ── Sheet 2: 기능 목록 ──
    ws2 = wb.create_sheet("기능 목록")
    headers2 = ["기능 영역", "기능명", "상세 설명", "관련 파일", "기술적 구현"]
    ws2.append(headers2)

    features = [
        ["3D 뷰어", "STL 파일 로딩", "ASCII & Binary STL 동시 지원, 다중 모델 로딩", "STLViewer.tsx, stl-loader.utils.ts", "BabylonJS STLLoader"],
        ["3D 뷰어", "카메라 컨트롤", "궤도(Orbit)/팬(Pan)/줌(Zoom) 카메라", "babylon.utils.ts", "ArcRotateCamera"],
        ["3D 뷰어", "Gizmo 조작", "마우스 드래그로 이동/회전/크기 변환", "babylon.utils.ts", "GizmoManager"],
        ["3D 뷰어", "가시성 토글", "파일별 표시/숨김 전환", "STLFileList.tsx", "mesh.isVisible"],
        ["3D 뷰어", "투명도 조절", "파일별 투명도 0~100% 조절", "STLFileList.tsx, STLViewer.tsx", "material.alpha"],
        ["3D 뷰어", "모델 하이라이트", "선택된 모델 강조 표시", "STLViewer.tsx", "커스텀 머티리얼 색상"],
        ["프로젝트", "프로젝트 CRUD", "생성/조회/수정/삭제", "ProjectPage.tsx, project.service.ts", "REST API + SQLite"],
        ["프로젝트", "환자 정보 관리", "환자 이름/코드 등 JSON 저장", "ProjectPage.tsx", "patientInfo JSON 필드"],
        ["프로젝트", "프로젝트 코드", "8자 랜덤 코드 자동 생성", "project.service.ts", "랜덤 문자열 생성"],
        ["파일 관리", "STL 파일 업로드", "파일 선택 다이얼로그로 업로드", "stl.service.ts, stl.controller.ts", "Multer 멀티파트"],
        ["파일 관리", "로컬 경로 임포트", "PC 파일시스템에서 직접 임포트", "LocalFileBrowser.tsx, fs.routes.ts", "fs.copyFile + Multer"],
        ["파일 관리", "파일시스템 탐색", "드라이브/폴더 직접 탐색", "LocalFileBrowser.tsx, fs.routes.ts", "Node.js fs.readdir"],
        ["변환", "Position 조절", "X/Y/Z 수치 입력 + Gizmo 드래그", "TransformPanel.tsx", "BabylonJS mesh.position"],
        ["변환", "Rotation 조절", "오일러각 입력 + Gizmo 드래그", "TransformPanel.tsx", "BabylonJS mesh.rotation"],
        ["변환", "Scale 조절", "균등/축별 스케일링", "TransformPanel.tsx", "BabylonJS mesh.scaling"],
        ["이력", "변환 이력 기록", "모든 변환 자동 로깅", "useAdjustmentLogs.ts", "adjustment_logs 테이블"],
        ["이력", "이력 조회", "STL 파일별 최신순 정렬", "HistoryViewer.tsx", "GET /api/stl/{id}/logs"],
        ["이력", "이력 삭제", "개별 또는 전체 삭제", "HistoryViewer.tsx", "DELETE API"],
        ["FDM 슬라이서", "메시 슬라이싱", "Z-plane 교차 알고리즘", "SliceEngine.ts", "삼각형-평면 교차 계산"],
        ["FDM 슬라이서", "벽 생성 (Walls)", "윤곽선 오프셋으로 동심원 벽", "GCodeGenerator.ts, PolygonClipper.ts", "js-clipper 오프셋"],
        ["FDM 슬라이서", "인필 채우기", "lines/grid/zigzag 패턴", "GCodeGenerator.ts", "라인 교차 알고리즘"],
        ["FDM 슬라이서", "G-code 생성", "G0/G1 명령어 + E 압출값", "GCodeGenerator.ts", "RepRap G-code 표준"],
        ["DLP 슬라이서", "마스크 이미지 생성", "흰색(노출)/검정(비노출) PNG", "ImageGenerator.ts", "OffscreenCanvas"],
        ["DLP 슬라이서", "프로젝터 해상도 설정", "X×Y 픽셀 + 픽셀 크기 설정", "SlicerPanel.tsx, types.ts", "DLPSettings 타입"],
        ["DLP 슬라이서", "노출 시간 설정", "레이어별 광원 노출 시간", "SlicerPanel.tsx", "exposureTime 파라미터"],
        ["인프라", "Web Worker 슬라이싱", "UI 비차단 병렬 처리", "SlicerWorker.ts, SlicerService.ts", "Web Worker API"],
        ["인프라", "로컬 전용 서버", "인터넷 불필요, localhost만 사용", "server.config.ts, index.ts", "Express 단일 서버"],
    ]
    for row in features:
        ws2.append(row)

    ws2.column_dimensions['A'].width = 14
    ws2.column_dimensions['B'].width = 20
    ws2.column_dimensions['C'].width = 40
    ws2.column_dimensions['D'].width = 35
    ws2.column_dimensions['E'].width = 28
    style_header(ws2)
    style_body(ws2)

    # ── Sheet 3: API 엔드포인트 ──
    ws3 = wb.create_sheet("API 엔드포인트")
    headers3 = ["메서드", "엔드포인트", "설명", "요청 본문", "응답", "관련 파일"]
    ws3.append(headers3)

    apis = [
        ["GET", "/health", "서버 상태 확인", "-", "{status, timestamp, environment}", "index.ts"],
        ["GET", "/api/fs?path=...", "로컬 파일시스템 탐색", "-", "[{name, type, size}]", "fs.routes.ts"],
        ["GET", "/api/projects?ownerId=...", "프로젝트 목록 조회", "-", "[Project]", "project.controller.ts"],
        ["GET", "/api/projects/{id}", "단일 프로젝트 조회", "-", "Project", "project.controller.ts"],
        ["POST", "/api/projects", "프로젝트 생성", "{ownerId, projectName, patientInfo?}", "Project", "project.controller.ts"],
        ["PUT", "/api/projects/{id}", "프로젝트 수정", "{projectName?, patientInfo?}", "Project", "project.controller.ts"],
        ["DELETE", "/api/projects/{id}", "프로젝트 삭제", "-", "{success}", "project.controller.ts"],
        ["GET", "/api/stl?projectId=...", "프로젝트 STL 파일 목록", "-", "[STLFile]", "stl.controller.ts"],
        ["POST", "/api/stl/upload", "STL 파일 업로드", "multipart/form-data (file, projectId)", "STLFile", "stl.controller.ts"],
        ["POST", "/api/stl/import-path", "로컬 경로 임포트", "{projectId, filePath}", "STLFile", "stl.controller.ts"],
        ["PUT", "/api/stl/{id}/visibility", "가시성 토글", "{visibility}", "STLFile", "stl.controller.ts"],
        ["PUT", "/api/stl/{id}/transform", "변환값 업데이트", "{transform}", "STLFile", "stl.controller.ts"],
        ["DELETE", "/api/stl/{id}", "STL 파일 삭제", "-", "{success}", "stl.controller.ts"],
        ["GET", "/api/stl/{id}/logs", "조정 이력 조회", "-", "[AdjustmentLog]", "stl.controller.ts"],
        ["POST", "/api/stl/{id}/logs", "조정 이력 생성", "{adjustmentType, deltaValue}", "AdjustmentLog", "stl.controller.ts"],
        ["DELETE", "/api/stl/{id}/logs/{logId}", "개별 이력 삭제", "-", "{success}", "stl.controller.ts"],
        ["DELETE", "/api/stl/{id}/logs", "전체 이력 초기화", "-", "{success}", "stl.controller.ts"],
    ]
    for row in apis:
        ws3.append(row)

    ws3.column_dimensions['A'].width = 10
    ws3.column_dimensions['B'].width = 30
    ws3.column_dimensions['C'].width = 22
    ws3.column_dimensions['D'].width = 35
    ws3.column_dimensions['E'].width = 30
    ws3.column_dimensions['F'].width = 22
    style_header(ws3)
    style_body(ws3)

    # ── Sheet 4: 참고 기술/라이브러리 분석 ──
    ws4 = wb.create_sheet("참고 기술 분석")
    headers4 = ["영역", "참고 대상", "유형", "MazicAlign 적용 방식", "근거/단서"]
    ws4.append(headers4)

    refs = [
        ["3D 엔진", "BabylonJS 공식 문서", "라이브러리 문서", "STL Loader, GizmoManager, ArcRotateCamera 활용", "코드에서 BabylonJS 6.x API 정확히 사용"],
        ["3D 엔진", "BabylonJS Playground", "예제 코드", "Scene/Camera/Light 설정 패턴", "babylon.utils.ts에 표준 설정 패턴 동일"],
        ["슬라이서", "Cura / PrusaSlicer", "오픈소스 참고", "Z-plane 교차, 벽 생성, 인필 패턴 알고리즘 참고", "SliceEngine.ts 알고리즘이 표준 FDM 슬라이싱 방식 준수"],
        ["슬라이서", "RepRap G-code 표준", "표준 명세", "G0/G1 + E값 G-code 생성", "GCodeGenerator.ts 출력이 RepRap 표준 준수"],
        ["슬라이서", "ChiTuBox / Lychee", "상용 참고", "DLP 마스크 이미지 생성 방식", "ImageGenerator.ts의 흰/검 마스크 방식이 표준 DLP 워크플로우"],
        ["폴리곤 연산", "Angus Johnson Clipper Library", "알고리즘 참고", "js-clipper로 폴리곤 오프셋 & Boolean 연산", "PolygonClipper.ts에서 js-clipper 직접 사용"],
        ["UI/UX", "3Shape / exocad", "상용 참고 가능성", "치과 3D 뷰어 UI 패턴 (좌측 파일목록, 중앙 뷰어, 우측 패널)", "치과 전문 뷰어의 전형적 3-패널 레이아웃"],
        ["UI/UX", "Blender / Meshmixer", "3D 도구 참고", "Transform Panel, Gizmo 조작 UX", "TransformPanel.tsx의 X/Y/Z 수치 + Gizmo 패턴"],
        ["프레임워크", "React 공식 문서", "프레임워크 문서", "함수형 컴포넌트 + 커스텀 훅 패턴", "모든 컴포넌트가 FC, 커스텀 훅(use*) 사용"],
        ["프레임워크", "Vite 공식 문서", "빌드 도구 문서", "React 플러그인 + API 프록시 설정", "vite.config.ts 설정이 Vite 표준 패턴"],
        ["프레임워크", "Zustand 공식 문서", "상태 관리 문서", "create() 스토어 패턴", "경량 상태 관리로 Zustand 선택"],
        ["백엔드", "Express 공식 가이드", "프레임워크 문서", "Controller → Service → DB 레이어 패턴", "표준 MVC 패턴 준수"],
        ["백엔드", "better-sqlite3 문서", "라이브러리 문서", "동기식 SQLite API + WAL 모드", "database.ts에 WAL 모드 설정 확인"],
        ["백엔드", "Multer 공식 예제", "라이브러리 문서", "멀티파트 파일 업로드 처리", "stl.controller.ts에서 Multer 미들웨어 사용"],
        ["V2 계획", "Three.js + react-three/fiber", "차기 기술", "BabylonJS 대체 계획 (번들 크기 감소)", "DESIGN.md에 명시적 기술 전환 로드맵"],
        ["V2 계획", "@react-three/drei", "차기 기술", "TransformControls 대체", "DESIGN.md에 Gizmo 대체 계획"],
        ["레거시", "Firebase", "레거시 흔적", "firebase.config.ts 파일 존재하나 미사용", "초기 Firebase 기반 → 로컬 SQLite로 전환한 흔적"],
    ]
    for row in refs:
        ws4.append(row)

    ws4.column_dimensions['A'].width = 12
    ws4.column_dimensions['B'].width = 28
    ws4.column_dimensions['C'].width = 16
    ws4.column_dimensions['D'].width = 40
    ws4.column_dimensions['E'].width = 45
    style_header(ws4)
    style_body(ws4)

    # ── Sheet 5: 파일 구조 ──
    ws5 = wb.create_sheet("파일 구조")
    headers5 = ["경로", "영역", "역할", "주요 내용"]
    ws5.append(headers5)

    files = [
        ["frontend/src/App.tsx", "Frontend", "라우터 설정", "React Router (/ → /projects, /viewer/:id)"],
        ["frontend/src/main.tsx", "Frontend", "엔트리 포인트", "React DOM mount"],
        ["frontend/src/pages/ProjectPage.tsx", "Frontend", "프로젝트 목록 페이지", "프로젝트 CRUD UI"],
        ["frontend/src/pages/ViewerPage.tsx", "Frontend", "3D 뷰어 메인 페이지", "뷰어/패널/슬라이서 통합 UI"],
        ["frontend/src/components/STLViewer.tsx", "Frontend", "3D 렌더러", "BabylonJS 캔버스 래퍼"],
        ["frontend/src/components/STLFileList.tsx", "Frontend", "파일 목록", "가시성/삭제/선택 UI"],
        ["frontend/src/components/LocalFileBrowser.tsx", "Frontend", "파일 탐색기", "PC 파일시스템 탐색 UI"],
        ["frontend/src/components/TransformPanel.tsx", "Frontend", "변환 패널", "Position/Rotation/Scale 입력"],
        ["frontend/src/components/HistoryViewer.tsx", "Frontend", "이력 뷰어", "조정 로그 목록"],
        ["frontend/src/components/Slicer/SlicerPanel.tsx", "Frontend", "슬라이서 UI", "FDM/DLP 설정 탭"],
        ["frontend/src/components/Slicer/SlicePreview.tsx", "Frontend", "슬라이스 미리보기", "레이어별 프리뷰"],
        ["frontend/src/services/slicer/SliceEngine.ts", "Frontend", "슬라이싱 코어", "Z-plane 교차 알고리즘"],
        ["frontend/src/services/slicer/GCodeGenerator.ts", "Frontend", "G-code 생성기", "FDM 벽/인필/경로 생성"],
        ["frontend/src/services/slicer/ImageGenerator.ts", "Frontend", "마스크 생성기", "DLP 흰/검 PNG 생성"],
        ["frontend/src/services/slicer/PolygonClipper.ts", "Frontend", "폴리곤 연산", "js-clipper 래퍼"],
        ["frontend/src/services/slicer/SlicerService.ts", "Frontend", "슬라이서 조정기", "Worker 관리 & 진행률"],
        ["frontend/src/services/slicer/SlicerWorker.ts", "Frontend", "Web Worker", "비차단 슬라이싱 실행"],
        ["frontend/src/utils/babylon.utils.ts", "Frontend", "BabylonJS 유틸", "Scene/Camera/Light/Gizmo 설정"],
        ["frontend/src/utils/stl-loader.utils.ts", "Frontend", "STL 로더 유틸", "STL 로딩 & 변환 적용"],
        ["frontend/src/hooks/useAuth.ts", "Frontend", "인증 훅", "로컬 사용자 상태"],
        ["frontend/src/hooks/useProjects.ts", "Frontend", "프로젝트 훅", "프로젝트 CRUD 로직"],
        ["frontend/src/hooks/useSTLFiles.ts", "Frontend", "STL 파일 훅", "파일 관리 로직"],
        ["frontend/src/hooks/useAdjustmentLogs.ts", "Frontend", "이력 훅", "조정 로그 관리"],
        ["backend/src/index.ts", "Backend", "서버 엔트리", "Express 앱 설정"],
        ["backend/src/config/database.ts", "Backend", "DB 초기화", "SQLite 스키마 생성"],
        ["backend/src/config/server.config.ts", "Backend", "서버 설정", "PORT, NODE_ENV 설정"],
        ["backend/src/controllers/project.controller.ts", "Backend", "프로젝트 컨트롤러", "HTTP 요청 핸들러"],
        ["backend/src/controllers/stl.controller.ts", "Backend", "STL 컨트롤러", "HTTP 요청 핸들러"],
        ["backend/src/services/project.service.ts", "Backend", "프로젝트 서비스", "DB 비즈니스 로직"],
        ["backend/src/services/stl.service.ts", "Backend", "STL 서비스", "DB 비즈니스 로직"],
        ["backend/src/routes/fs.routes.ts", "Backend", "파일시스템 라우트", "파일 탐색 API"],
        ["backend/src/routes/project.routes.ts", "Backend", "프로젝트 라우트", "프로젝트 CRUD 엔드포인트"],
        ["backend/src/routes/stl.routes.ts", "Backend", "STL 라우트", "STL CRUD 엔드포인트"],
    ]
    for row in files:
        ws5.append(row)

    ws5.column_dimensions['A'].width = 48
    ws5.column_dimensions['B'].width = 12
    ws5.column_dimensions['C'].width = 18
    ws5.column_dimensions['D'].width = 40
    style_header(ws5)
    style_body(ws5)

    # ── Sheet 6: PrusaSlicer 기능 비교 ──
    ws6 = wb.create_sheet("PrusaSlicer 기능 비교")
    headers6 = ["기능 영역", "기능명", "PrusaSlicer", "MazicAlign", "Gap 수준",
                "우선순위", "구현 난이도", "PrusaSlicer 참고 파일/모듈", "비고"]
    ws6.append(headers6)

    gap_fill = PatternFill(start_color="FEE2E2", end_color="FEE2E2", fill_type="solid")  # red tint
    ok_fill = PatternFill(start_color="D1FAE5", end_color="D1FAE5", fill_type="solid")   # green tint

    cmp_rows = [
        # [영역, 기능, PS, MA, Gap, 우선순위, 난이도, 참고파일, 비고]
        ["인필 패턴", "Lines / Grid / Zigzag", "O", "O", "동등", "-", "-", "Fill/FillLine.cpp, FillRectilinear.cpp", "MazicAlign 기본 구현 완료"],
        ["인필 패턴", "Honeycomb (벌집)", "O", "X", "부족", "높음", "중", "Fill/FillHoneycomb.cpp", "치과 모델 내부 강도에 유용"],
        ["인필 패턴", "3D Honeycomb", "O", "X", "부족", "중간", "높음", "Fill/Fill3DHoneycomb.cpp", "3D 벌집 — 등방성 강도"],
        ["인필 패턴", "Gyroid (자이로이드)", "O", "X", "부족", "높음", "중", "Fill/FillGyroid.cpp", "레진 배출 용이, SLA/DLP에 적합"],
        ["인필 패턴", "Concentric (동심원)", "O", "X", "부족", "높음", "낮음", "Fill/FillConcentric.cpp", "상/하면에 적합, 구현 간단"],
        ["인필 패턴", "Lightning (번개)", "O", "X", "부족", "낮음", "높음", "Fill/FillLightning.cpp, Lightning/", "재료 절약형, 트리 알고리즘"],
        ["인필 패턴", "Adaptive (적응형)", "O", "X", "부족", "중간", "높음", "Fill/FillAdaptive.cpp", "모델 형상에 따라 밀도 자동 조절"],
        ["인필 패턴", "PlanePath 패턴", "O", "X", "부족", "낮음", "중", "Fill/FillPlanePath.cpp", "Hilbert 곡선 등 수학적 패턴"],
        ["벽 생성", "기본 Walls/Perimeters", "O", "O", "동등", "-", "-", "GCode/ExtrusionOrder.cpp", "MazicAlign js-clipper 오프셋 사용"],
        ["벽 생성", "Arachne 가변폭 윤곽선", "O", "X", "부족", "중간", "매우 높음", "Arachne/", "가변 폭 압출 — 얇은 벽 품질 향상"],
        ["벽 생성", "Seam 배치 최적화", "O", "X", "부족", "높음", "중", "GCode/Seam*.cpp (10+ 전략)", "이음새 위치 최적화로 외관 품질 향상"],
        ["벽 생성", "Gap Fill (틈 채우기)", "O", "△", "부분", "중간", "중", "GCode/ExtrusionProcessor.cpp", "MA에 옵션 있으나 알고리즘 단순"],
        ["서포트", "일반 서포트 (그리드/라인)", "O", "X", "부족", "매우 높음", "높음", "Support/SupportMaterial.cpp", "FDM 필수 기능, 오버행 지지"],
        ["서포트", "트리 서포트 (Tree)", "O", "X", "부족", "높음", "매우 높음", "Support/TreeSupport.cpp", "재료 절약 + 표면 품질"],
        ["서포트", "유기적 서포트 (Organic)", "O", "X", "부족", "중간", "매우 높음", "Support/OrganicSupport.cpp", "곡선형 트리 서포트"],
        ["서포트", "서포트 파라미터 설정", "O", "X", "부족", "높음", "중", "Support/SupportParameters.cpp", "간격/밀도/패턴/접촉면적 등"],
        ["G-code", "기본 G0/G1 생성", "O", "O", "동등", "-", "-", "GCode/GCodeWriter.cpp", "MazicAlign GCodeGenerator.ts"],
        ["G-code", "리트랙션 (Retraction)", "O", "X", "부족", "매우 높음", "낮음", "GCode/GCodeWriter.cpp", "문자열화 방지 핵심 기능"],
        ["G-code", "압력 균등화", "O", "X", "부족", "중간", "중", "GCode/PressureEqualizer.cpp", "Linear Advance 등"],
        ["G-code", "냉각 제어", "O", "X", "부족", "높음", "중", "GCode/CoolingBuffer.cpp", "레이어별 팬 속도/대기 시간"],
        ["G-code", "다중 펌웨어 지원", "O", "X", "부족", "중간", "낮음", "GCode/GCodeWriter.cpp", "Marlin/Prusa/RepRap/Mach3 등"],
        ["G-code", "Wipe Tower (다중 재료)", "O", "X", "해당없음", "-", "-", "GCode/WipeTower*.cpp", "다중 익스트루더 전용"],
        ["G-code", "나선형 화병 모드", "O", "X", "부족", "중간", "중", "GCode/SpiralVase.cpp", "단일벽 연속 나선 출력"],
        ["G-code", "이동 경로 최적화", "O", "X", "부족", "높음", "중", "GCode/Travels.cpp", "출력 시간 단축"],
        ["G-code", "경계선 교차 회피", "O", "X", "부족", "중간", "높음", "GCode/AvoidCrossingPerimeters.cpp", "이동시 외벽 교차 방지"],
        ["G-code", "충돌 감지", "O", "X", "부족", "낮음", "높음", "GCode/ConflictChecker.cpp", "다중 모델 충돌 방지"],
        ["G-code", "후처리 스크립트", "O", "X", "부족", "중간", "낮음", "GCode/PostProcessor.cpp", "사용자 스크립트 실행"],
        ["G-code", "썸네일 생성", "O", "X", "부족", "낮음", "낮음", "GCode/Thumbnails.cpp", "프린터 LCD 미리보기"],
        ["메시 처리", "STL 자동 복구", "O", "X", "부족", "높음", "높음", "TriangleMesh.cpp", "깨진 메시 자동 수정"],
        ["메시 처리", "다중 포맷 (OBJ/AMF/3MF)", "O", "△ (STL만)", "부분", "높음", "중", "Format/", "OBJ 추가만 해도 실용성 증가"],
        ["메시 처리", "CSG Boolean 연산", "O", "X", "부족", "낮음", "매우 높음", "CSGMesh/", "메시 합집합/차집합/교집합"],
        ["메시 처리", "모델 자동 배치 (Arrange)", "O", "X", "부족", "중간", "높음", "slic3r-arrange/", "빌드 플레이트 최적 배치"],
        ["SLA/DLP", "마스크 이미지 생성", "O", "O", "동등", "-", "-", "SLA/RasterBase.cpp, AGGRaster.hpp", "MazicAlign ImageGenerator.ts"],
        ["SLA/DLP", "자동 서포트 포인트 생성", "O", "X", "부족", "매우 높음", "높음", "SLA/SupportPointGenerator.cpp", "SLA/DLP 핵심 기능"],
        ["SLA/DLP", "SLA 트리 서포트 구조", "O", "X", "부족", "매우 높음", "매우 높음", "SLA/SupportTree*.cpp", "포인트→트리 구조 변환"],
        ["SLA/DLP", "모델 중공화 (Hollowing)", "O", "X", "부족", "높음", "높음", "SLA/Hollowing.cpp", "레진 절약 & 내부 비움"],
        ["SLA/DLP", "빌드 패드 (Pad)", "O", "X", "부족", "높음", "중", "SLA/Pad.cpp", "빌드 플레이트 접착 보조"],
        ["SLA/DLP", "최적 인쇄 회전 (Rotfinder)", "O", "X", "부족", "중간", "높음", "SLA/Rotfinder.cpp", "서포트 최소화 회전각 탐색"],
        ["SLA/DLP", "Z축 보정 (ZCorrection)", "O", "X", "부족", "중간", "낮음", "SLA/ZCorrection.cpp", "레이어별 Z 오프셋 보정"],
        ["SLA/DLP", "래스터→폴리곤 변환", "O", "X", "부족", "낮음", "중", "SLA/RasterToPolygons.cpp", "래스터 이미지를 벡터로 변환"],
        ["기타", "CLI 모드 (GUI 없이)", "O", "X", "해당없음", "-", "-", "CLI/", "데스크톱 앱 전용 기능"],
        ["기타", "멀티스레드 처리", "O (C++ Thread)", "O (Web Worker)", "동등", "-", "-", "Execution/", "방식은 다르나 동일 목적"],
        ["기타", "유닛 테스트 자동화", "O", "X", "부족", "중간", "낮음", "tests/", "품질 보증 체계"],
        ["기타", "Clipper 폴리곤 라이브러리", "O (clipper/)", "O (js-clipper)", "동등", "-", "-", "clipper/", "동일 알고리즘의 C++/JS 버전"],
    ]
    for row in cmp_rows:
        ws6.append(row)

    ws6.column_dimensions['A'].width = 12
    ws6.column_dimensions['B'].width = 28
    ws6.column_dimensions['C'].width = 14
    ws6.column_dimensions['D'].width = 14
    ws6.column_dimensions['E'].width = 10
    ws6.column_dimensions['F'].width = 12
    ws6.column_dimensions['G'].width = 12
    ws6.column_dimensions['H'].width = 40
    ws6.column_dimensions['I'].width = 35
    style_header(ws6)
    style_body(ws6, max_col=9)

    # Color-code the Gap column (column E = 5)
    for ri in range(2, ws6.max_row + 1):
        gap_cell = ws6.cell(row=ri, column=5)
        if gap_cell.value == "부족" or gap_cell.value == "부분":
            gap_cell.fill = gap_fill
            gap_cell.font = Font(name="맑은 고딕", size=10, bold=True, color="DC2626")
        elif gap_cell.value == "동등":
            gap_cell.fill = ok_fill
            gap_cell.font = Font(name="맑은 고딕", size=10, bold=True, color="059669")

    # ── Sheet 7: PrusaSlicer 모듈 구조 ──
    ws7 = wb.create_sheet("PrusaSlicer 모듈 구조")
    headers7 = ["디렉토리", "모듈명", "역할", "주요 파일 수", "MazicAlign 대응", "참고 가치"]
    ws7.append(headers7)

    modules = [
        ["src/libslic3r/", "libslic3r (핵심 엔진)", "슬라이싱 엔진 라이브러리", "150+", "SliceEngine.ts + GCodeGenerator.ts", "매우 높음"],
        ["src/libslic3r/Fill/", "Fill (인필 패턴)", "10+ 인필 패턴 구현", "24", "GCodeGenerator.ts (3 패턴만)", "높음"],
        ["src/libslic3r/GCode/", "GCode (코드 생성)", "G-code 생성/최적화/후처리", "66", "GCodeGenerator.ts", "매우 높음"],
        ["src/libslic3r/GCode/Seam*", "Seam (이음새)", "10+ 이음새 배치 전략", "~20", "없음", "높음"],
        ["src/libslic3r/Support/", "Support (FDM 서포트)", "일반/트리/유기적 서포트", "17", "없음", "매우 높음"],
        ["src/libslic3r/SLA/", "SLA (광조형)", "SLA/DLP 전용 기능", "38", "ImageGenerator.ts (기본만)", "매우 높음"],
        ["src/libslic3r/Arachne/", "Arachne (가변폭)", "가변 폭 윤곽선 엔진", "-", "없음", "중간"],
        ["src/libslic3r/Algorithm/", "Algorithm (알고리즘)", "범용 알고리즘 유틸", "-", "없음", "중간"],
        ["src/libslic3r/Geometry/", "Geometry (기하학)", "기하학 연산 라이브러리", "-", "transform.utils.ts (기본만)", "높음"],
        ["src/libslic3r/Format/", "Format (파일 포맷)", "STL/OBJ/AMF/3MF 입출력", "-", "stl-loader.utils.ts (STL만)", "높음"],
        ["src/libslic3r/Optimize/", "Optimize (최적화)", "경로/배치 최적화", "-", "없음", "중간"],
        ["src/libslic3r/CSGMesh/", "CSGMesh (Boolean)", "메시 Boolean 연산", "-", "없음", "낮음"],
        ["src/libslic3r/BranchingTree/", "BranchingTree", "분기형 트리 구조 생성", "-", "없음", "중간"],
        ["src/clipper/", "Clipper (폴리곤)", "Angus Johnson Clipper C++", "-", "js-clipper (JS 포팅)", "동일 알고리즘"],
        ["src/slic3r/", "slic3r (GUI)", "wxWidgets GUI 애플리케이션", "-", "React UI (별도 접근)", "참고만"],
        ["src/slic3r-arrange/", "Arrange (자동 배치)", "빌드 플레이트 자동 배치", "-", "없음", "중간"],
        ["src/libvgcode/", "VGCode (벡터 G-code)", "벡터 G-code 시각화", "-", "SlicePreview.tsx (기본만)", "중간"],
        ["src/CLI/", "CLI (명령행)", "GUI 없이 슬라이싱", "-", "없음 (웹 전용)", "해당없음"],
    ]
    for row in modules:
        ws7.append(row)

    ws7.column_dimensions['A'].width = 30
    ws7.column_dimensions['B'].width = 25
    ws7.column_dimensions['C'].width = 30
    ws7.column_dimensions['D'].width = 12
    ws7.column_dimensions['E'].width = 32
    ws7.column_dimensions['F'].width = 14
    style_header(ws7)
    style_body(ws7)

    wb.save("c:/Users/JoWooHyun/Documents/MazicAlign/MazicAlign_분석정리표_v2.xlsx")
    print("XLSX 생성 완료: MazicAlign_분석정리표_v2.xlsx")


if __name__ == "__main__":
    create_pptx()
    create_xlsx()
